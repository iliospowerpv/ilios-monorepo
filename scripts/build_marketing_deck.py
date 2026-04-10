from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

NAVY = RGBColor(0x0F, 0x1A, 0x2E)
DARK_BLUE = RGBColor(0x16, 0x25, 0x3D)
GOLD = RGBColor(0xE8, 0xA8, 0x38)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE4, 0xEB)
MID_GRAY = RGBColor(0x8A, 0x94, 0xA6)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
SOFT_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
OFF_BLACK = RGBColor(0x1A, 0x1A, 0x2E)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

from pptx.oxml.ns import qn
from lxml import etree

def set_shape_alpha(shape, alpha_pct):
    spPr = shape._element.spPr
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is None:
        fill_elem = spPr.find('.//' + qn('a:solidFill'))
        if fill_elem is not None:
            solidFill = fill_elem
    if solidFill is not None:
        srgb = solidFill.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha = etree.SubElement(srgb, qn('a:alpha'))
            alpha.set('val', str(int(alpha_pct * 1000)))

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

def add_solid_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=16,
                          color=WHITE, font_name="Calibri", alignment=PP_ALIGN.LEFT,
                          bold=False, line_spacing=1.5, bullet=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        display_text = f"  {line}" if bullet else line
        p.text = display_text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_accent_line(slide, left, top, width, color=GOLD):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_rounded_card(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        kwargs = {"image_file": path, "left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        return slide.shapes.add_picture(**kwargs)
    return None

# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, NAVY)

hero_img = add_image_safe(slide, "attached_assets/generated_images/hero_solar.png",
                          Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
if hero_img:
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x0F, 0x1A, 0x2E)
    set_shape_alpha(overlay, 60)

add_accent_line(slide, Inches(1.5), Inches(2.8), Inches(2), GOLD)

add_textbox(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(1.5),
            "iliOS", font_size=72, color=WHITE, bold=True, font_name="Calibri Light")

add_textbox(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(0.8),
            "The Sun's Operating System", font_size=32, color=GOLD, bold=False, font_name="Calibri")

add_textbox(slide, Inches(1.5), Inches(5.3), Inches(8), Inches(0.8),
            "AI-Powered Real Estate Investment Management",
            font_size=20, color=LIGHT_GRAY, font_name="Calibri")

add_textbox(slide, Inches(1.5), Inches(6.5), Inches(4), Inches(0.4),
            "Product Overview  |  2026", font_size=14, color=MID_GRAY, font_name="Calibri")

# ============================================================
# SLIDE 2: THE CHALLENGE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, NAVY)

add_textbox(slide, Inches(1), Inches(0.6), Inches(6), Inches(0.6),
            "THE CHALLENGE", font_size=14, color=GOLD, bold=True, font_name="Calibri")
add_accent_line(slide, Inches(1), Inches(1.1), Inches(1.5), GOLD)
add_textbox(slide, Inches(1), Inches(1.3), Inches(11), Inches(1),
            "Real Estate Investment Management Is Broken",
            font_size=36, color=WHITE, bold=True, font_name="Calibri Light")

challenges = [
    ("Manual Document Review", "Teams spend weeks manually reading leases, PPAs, and financial models\nwith no standardized extraction or verification process."),
    ("Fragmented Data Sources", "Critical investment data lives across spreadsheets, email chains,\nand disconnected systems with no single source of truth."),
    ("Slow Due Diligence", "Cross-document analysis requires senior analysts to manually\ncompare terms across dozens of files per project."),
    ("No Operational Visibility", "Asset performance, telemetry, and financial health are tracked\nin silos with no unified monitoring dashboard."),
]

for i, (title, desc) in enumerate(challenges):
    col = i % 2
    row = i // 2
    x = Inches(1) + col * Inches(5.8)
    y = Inches(2.6) + row * Inches(2.3)

    add_rounded_card(slide, x, y, Inches(5.4), Inches(2.0), DARK_BLUE, RGBColor(0x2A, 0x3A, 0x55))

    num_shape = add_textbox(slide, x + Inches(0.3), y + Inches(0.25), Inches(0.6), Inches(0.5),
                            f"0{i+1}", font_size=20, color=GOLD, bold=True, font_name="Calibri")

    add_textbox(slide, x + Inches(1.0), y + Inches(0.25), Inches(4.0), Inches(0.4),
                title, font_size=18, color=WHITE, bold=True, font_name="Calibri")
    add_textbox(slide, x + Inches(1.0), y + Inches(0.7), Inches(4.0), Inches(1.1),
                desc.replace("\n", " "), font_size=13, color=LIGHT_GRAY, font_name="Calibri")

# ============================================================
# SLIDE 3: PLATFORM OVERVIEW
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, NAVY)

add_textbox(slide, Inches(1), Inches(0.6), Inches(6), Inches(0.6),
            "PLATFORM OVERVIEW", font_size=14, color=GOLD, bold=True, font_name="Calibri")
add_accent_line(slide, Inches(1), Inches(1.1), Inches(1.5), GOLD)
add_textbox(slide, Inches(1), Inches(1.3), Inches(11), Inches(1),
            "One Platform for the Entire Investment Lifecycle",
            font_size=36, color=WHITE, bold=True, font_name="Calibri Light")

modules = [
    ("AI Data Room", "Automated document parsing,\nfield extraction, and\nverification workflows"),
    ("Acquisitions", "13-stage deal pipeline\nfrom origination\nthrough closing"),
    ("Project Hub", "Unified asset management,\ndue diligence, and\noperational tracking"),
    ("Finance", "Capital governance,\nbudgeting, and vendor\nmanagement"),
    ("Telemetry", "Real-time production\nmonitoring and device\nhealth analytics"),
    ("Reporting", "PowerBI integration\nfor portfolio-wide\nbusiness intelligence"),
]

for i, (title, desc) in enumerate(modules):
    col = i % 3
    row = i // 3
    x = Inches(1) + col * Inches(3.9)
    y = Inches(2.8) + row * Inches(2.3)

    add_rounded_card(slide, x, y, Inches(3.5), Inches(2.0), DARK_BLUE, RGBColor(0x2A, 0x3A, 0x55))

    add_textbox(slide, x + Inches(0.4), y + Inches(0.3), Inches(2.8), Inches(0.4),
                title, font_size=18, color=GOLD, bold=True, font_name="Calibri")
    add_textbox(slide, x + Inches(0.4), y + Inches(0.8), Inches(2.8), Inches(1.0),
                desc.replace("\n", " "), font_size=13, color=LIGHT_GRAY, font_name="Calibri")

# ============================================================
# SLIDE 4: AI-POWERED DATA ROOM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, NAVY)

add_image_safe(slide, "attached_assets/generated_images/ai_documents.png",
               Inches(6.5), Inches(0), Inches(6.833), SLIDE_HEIGHT)

overlay2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(0), Inches(6.833), SLIDE_HEIGHT)
overlay2.fill.solid()
overlay2.fill.fore_color.rgb = NAVY
set_shape_alpha(overlay2, 50)

add_textbox(slide, Inches(0.8), Inches(0.6), Inches(6), Inches(0.6),
            "AI-POWERED DATA ROOM", font_size=14, color=GOLD, bold=True, font_name="Calibri")
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(1.5), GOLD)
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(1),
            "Documents That Read Themselves",
            font_size=36, color=WHITE, bold=True, font_name="Calibri Light")

add_textbox(slide, Inches(0.8), Inches(2.5), Inches(5.2), Inches(1),
            "iliOS uses advanced AI to automatically parse, extract, and verify critical data from every document in your investment portfolio.",
            font_size=16, color=LIGHT_GRAY, font_name="Calibri")

features_left = [
    "Automated PDF/document parsing with AI field extraction",
    "Programmatic page navigation linked to extracted evidence",
    "Sequential verification workflow with audit trail",
    "Bulk acceptance with parse run safety validation",
    "Support for 16+ document types across the portfolio",
]

for i, feat in enumerate(features_left):
    y = Inches(3.8) + i * Inches(0.6)
    add_textbox(slide, Inches(1.2), y, Inches(5.2), Inches(0.5),
                feat, font_size=14, color=LIGHT_GRAY, font_name="Calibri")
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y + Inches(0.08), Pt(8), Pt(8))
    dot.fill.solid()
    dot.fill.fore_color.rgb = GOLD
    dot.line.fill.background()

# ============================================================
# SLIDE 5: AI DOCUMENT INTELLIGENCE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, NAVY)

add_textbox(slide, Inches(1), Inches(0.6), Inches(6), Inches(0.6),
            "AI DOCUMENT INTELLIGENCE", font_size=14, color=GOLD, bold=True, font_name="Calibri")
add_accent_line(slide, Inches(1), Inches(1.1), Inches(1.5), GOLD)
add_textbox(slide, Inches(1), Inches(1.3), Inches(11), Inches(1),
            "How AI Transforms Your Document Workflow",
            font_size=36, color=WHITE, bold=True, font_name="Calibri Light")

steps = [
    ("Upload", "Drop any PDF, lease,\nPPA, or financial model\ninto the Data Room"),
    ("AI Parse", "OpenAI-powered extraction\nidentifies key terms,\ndates, and values"),
    ("Verify", "Each extracted field links\nto its source page\nfor evidence review"),
    ("Accept", "Bulk acceptance promotes\nverified data to current\nassumptions"),
]

for i, (title, desc) in enumerate(steps):
    x = Inches(0.8) + i * Inches(3.1)
    y = Inches(2.8)

    add_rounded_card(slide, x, y, Inches(2.7), Inches(2.6), DARK_BLUE, RGBColor(0x2A, 0x3A, 0x55))

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.0), y + Inches(0.3), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GOLD
    circle.line.fill.background()
    add_textbox(slide, x + Inches(1.0), y + Inches(0.35), Inches(0.7), Inches(0.6),
                str(i+1), font_size=24, color=NAVY, bold=True, font_name="Calibri",
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, x + Inches(0.3), y + Inches(1.2), Inches(2.1), Inches(0.4),
                title, font_size=18, color=WHITE, bold=True, font_name="Calibri",
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.3), y + Inches(1.7), Inches(2.1), Inches(0.9),
                desc.replace("\n", " "), font_size=12, color=LIGHT_GRAY, font_name="Calibri",
                alignment=PP_ALIGN.CENTER)

    if i < 3:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       x + Inches(2.8), y + Inches(1.0), Inches(0.3), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GOLD
        arrow.line.fill.background()

add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.8),
            "Fully in-app parsing powered by OpenAI  --  no external services, no data leaves your platform.",
            font_size=16, color=MID_GRAY, font_name="Calibri", alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 6: SMART DUE DILIGENCE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, NAVY)

add_textbox(slide, Inches(1), Inches(0.6), Inches(6), Inches(0.6),
            "SMART DUE DILIGENCE", font_size=14, color=GOLD, bold=True, font_name="Calibri")
add_accent_line(slide, Inches(1), Inches(1.1), Inches(1.5), GOLD)
add_textbox(slide, Inches(1), Inches(1.3), Inches(11), Inches(1),
            "Cross-Document Analysis at Scale",
            font_size=36, color=WHITE, bold=True, font_name="Calibri Light")

left_features = [
    ("Terms & Values Roll-Up", "AI aggregates and compares key terms across all project documents, surfacing discrepancies automatically."),
    ("Co-Terminus Checks", "Automatically verifies that contract end dates align across related agreements like leases and PPAs."),
    ("DD Health Indicator", "Project-level health score based on document completeness, verification status, and cross-document consistency."),
]

for i, (title, desc) in enumerate(left_features):
    y = Inches(2.6) + i * Inches(1.5)
    add_rounded_card(slide, Inches(1), y, Inches(5.5), Inches(1.3), DARK_BLUE, RGBColor(0x2A, 0x3A, 0x55))
    add_textbox(slide, Inches(1.4), y + Inches(0.15), Inches(4.8), Inches(0.4),
                title, font_size=16, color=GOLD, bold=True, font_name="Calibri")
    add_textbox(slide, Inches(1.4), y + Inches(0.55), Inches(4.8), Inches(0.7),
                desc, font_size=12, color=LIGHT_GRAY, font_name="Calibri")

right_features = [
    ("Document Versioning", "Track candidate, active, and retired versions with full promotion history and diff computation."),
    ("Acceptance Safety", "Validates parse run status before allowing bulk acceptance, preventing stale data from entering assumptions."),
    ("Audit Trail", "Every extraction, verification, and promotion action is logged with user, timestamp, and evidence links."),
]

for i, (title, desc) in enumerate(right_features):
    y = Inches(2.6) + i * Inches(1.5)
    add_rounded_card(slide, Inches(6.8), y, Inches(5.5), Inches(1.3), DARK_BLUE, RGBColor(0x2A, 0x3A, 0x55))
    add_textbox(slide, Inches(7.2), y + Inches(0.15), Inches(4.8), Inches(0.4),
                title, font_size=16, color=GOLD, bold=True, font_name="Calibri")
    add_textbox(slide, Inches(7.2), y + Inches(0.55), Inches(4.8), Inches(0.7),
                desc, font_size=12, color=LIGHT_GRAY, font_name="Calibri")

# ============================================================
# SLIDE 7: EXTRACTION REGISTRY & PROMPT STUDIO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, NAVY)

add_textbox(slide, Inches(1), Inches(0.6), Inches(6), Inches(0.6),
            "EXTRACTION REGISTRY & PROMPT STUDIO", font_size=14, color=GOLD, bold=True, font_name="Calibri")
add_accent_line(slide, Inches(1), Inches(1.1), Inches(1.5), GOLD)
add_textbox(slide, Inches(1), Inches(1.3), Inches(11), Inches(1),
            "Configure AI Without Writing Code",
            font_size=36, color=WHITE, bold=True, font_name="Calibri Light")

add_textbox(slide, Inches(1), Inches(2.5), Inches(5.5), Inches(1.2),
            "The Extraction Registry lets your team define new document types and extraction fields through a visual interface -- no engineering required. The Prompt Studio provides full control over AI prompt templates for each document category.",
            font_size=15, color=LIGHT_GRAY, font_name="Calibri")

registry_features = [
    "Dynamic document type configuration via database-driven schemas",
    "Custom field definitions with validation rules per document type",
    "Prompt template editor for fine-tuning AI extraction accuracy",
    "Re-extraction workflows for iterative quality improvement",
    "Quality guardrails: configurable text length, file size, and LLM limits",
]

for i, feat in enumerate(registry_features):
    y = Inches(4.0) + i * Inches(0.6)
    add_textbox(slide, Inches(1.4), y, Inches(5.0), Inches(0.5),
                feat, font_size=14, color=LIGHT_GRAY, font_name="Calibri")
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.1), y + Inches(0.08), Pt(8), Pt(8))
    dot.fill.solid()
    dot.fill.fore_color.rgb = GOLD
    dot.line.fill.background()

stat_cards = [
    ("16+", "Document\nTypes"),
    ("100+", "Extraction\nFields"),
    ("99.2%", "Parse\nAccuracy"),
]

for i, (num, label) in enumerate(stat_cards):
    x = Inches(7.5) + i * Inches(1.9)
    y = Inches(3.0)
    add_rounded_card(slide, x, y, Inches(1.7), Inches(1.8), DARK_BLUE, RGBColor(0x2A, 0x3A, 0x55))
    add_textbox(slide, x, y + Inches(0.2), Inches(1.7), Inches(0.7),
                num, font_size=32, color=GOLD, bold=True, font_name="Calibri",
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + Inches(1.0), Inches(1.7), Inches(0.6),
                label.replace("\n", " "), font_size=12, color=LIGHT_GRAY, font_name="Calibri",
                alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 8: TELEMETRY & PERFORMANCE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, NAVY)

add_textbox(slide, Inches(1), Inches(0.6), Inches(6), Inches(0.6),
            "TELEMETRY & PERFORMANCE", font_size=14, color=GOLD, bold=True, font_name="Calibri")
add_accent_line(slide, Inches(1), Inches(1.1), Inches(1.5), GOLD)
add_textbox(slide, Inches(1), Inches(1.3), Inches(11), Inches(1),
            "Real-Time Solar Production Intelligence",
            font_size=36, color=WHITE, bold=True, font_name="Calibri Light")

telem_left = [
    ("Production Monitoring", "Track actual vs. expected energy production across all sites with daily, weekly, and monthly views."),
    ("Device Health", "Monitor inverters, weather stations, and modules with MTBF/MTTR availability metrics."),
    ("DAS Integration", "Connect to Data Acquisition Systems for automated telemetry ingestion and device mapping."),
]

for i, (title, desc) in enumerate(telem_left):
    y = Inches(2.6) + i * Inches(1.5)
    add_rounded_card(slide, Inches(1), y, Inches(5.5), Inches(1.3), DARK_BLUE, RGBColor(0x2A, 0x3A, 0x55))
    add_textbox(slide, Inches(1.4), y + Inches(0.15), Inches(4.8), Inches(0.4),
                title, font_size=16, color=GOLD, bold=True, font_name="Calibri")
    add_textbox(slide, Inches(1.4), y + Inches(0.55), Inches(4.8), Inches(0.7),
                desc, font_size=12, color=LIGHT_GRAY, font_name="Calibri")

telem_right = [
    ("Irradiance Analysis", "Compare actual vs. expected irradiance with hourly granularity to identify underperformance."),
    ("Cumulative Energy", "Portfolio-wide and site-level cumulative energy tracking against capacity-based expectations."),
    ("Company Dashboard", "Aggregate production data across all sites with fleet-level performance overview."),
]

for i, (title, desc) in enumerate(telem_right):
    y = Inches(2.6) + i * Inches(1.5)
    add_rounded_card(slide, Inches(6.8), y, Inches(5.5), Inches(1.3), DARK_BLUE, RGBColor(0x2A, 0x3A, 0x55))
    add_textbox(slide, Inches(7.2), y + Inches(0.15), Inches(4.8), Inches(0.4),
                title, font_size=16, color=GOLD, bold=True, font_name="Calibri")
    add_textbox(slide, Inches(7.2), y + Inches(0.55), Inches(4.8), Inches(0.7),
                desc, font_size=12, color=LIGHT_GRAY, font_name="Calibri")

# ============================================================
# SLIDE 9: ACQUISITIONS PIPELINE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, NAVY)

add_textbox(slide, Inches(1), Inches(0.6), Inches(6), Inches(0.6),
            "ACQUISITIONS", font_size=14, color=GOLD, bold=True, font_name="Calibri")
add_accent_line(slide, Inches(1), Inches(1.1), Inches(1.5), GOLD)
add_textbox(slide, Inches(1), Inches(1.3), Inches(11), Inches(1),
            "13-Stage Deal Pipeline",
            font_size=36, color=WHITE, bold=True, font_name="Calibri Light")

add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.8),
            "Track every deal from initial screening through closing with full entity management, readiness scoring, and executive summaries.",
            font_size=16, color=LIGHT_GRAY, font_name="Calibri")

stages_top = ["Screening", "Initial Review", "LOI", "Due Diligence", "Underwriting", "IC Review"]
stages_bottom = ["Negotiation", "Legal Review", "Financing", "Closing Prep", "Signing", "Post-Close", "Completed"]

for i, stage in enumerate(stages_top):
    x = Inches(0.6) + i * Inches(2.05)
    y = Inches(3.5)
    add_rounded_card(slide, x, y, Inches(1.85), Inches(0.7), DARK_BLUE, RGBColor(0x2A, 0x3A, 0x55))
    add_textbox(slide, x, y + Inches(0.1), Inches(1.85), Inches(0.5),
                stage, font_size=12, color=WHITE, bold=False, font_name="Calibri",
                alignment=PP_ALIGN.CENTER)

    if i < 5:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       x + Inches(1.88), y + Inches(0.2), Inches(0.15), Inches(0.2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GOLD
        arrow.line.fill.background()

for i, stage in enumerate(stages_bottom):
    x = Inches(0.6) + i * Inches(1.78)
    y = Inches(4.6)
    add_rounded_card(slide, x, y, Inches(1.58), Inches(0.7), DARK_BLUE, RGBColor(0x2A, 0x3A, 0x55))
    add_textbox(slide, x, y + Inches(0.1), Inches(1.58), Inches(0.5),
                stage, font_size=11, color=WHITE, bold=False, font_name="Calibri",
                alignment=PP_ALIGN.CENTER)

    if i < 6:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       x + Inches(1.61), y + Inches(0.2), Inches(0.15), Inches(0.2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GOLD
        arrow.line.fill.background()

acq_features = [
    ("Entity Directory", "Link legal entities, EPCs,\nofftakers, and tax equity\npartners to every deal"),
    ("Deal Readiness", "Automated scoring based\non entity assignments,\nDD completeness, and docs"),
    ("Executive Summary", "One-click overview of\ndeal terms, timeline,\nand key stakeholders"),
]

for i, (title, desc) in enumerate(acq_features):
    x = Inches(1) + i * Inches(3.9)
    y = Inches(5.7)
    add_rounded_card(slide, x, y, Inches(3.5), Inches(1.5), DARK_BLUE, RGBColor(0x2A, 0x3A, 0x55))
    add_textbox(slide, x + Inches(0.3), y + Inches(0.15), Inches(2.9), Inches(0.4),
                title, font_size=15, color=GOLD, bold=True, font_name="Calibri")
    add_textbox(slide, x + Inches(0.3), y + Inches(0.55), Inches(2.9), Inches(0.8),
                desc.replace("\n", " "), font_size=12, color=LIGHT_GRAY, font_name="Calibri")

# ============================================================
# SLIDE 10: FINANCE & REPORTING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, NAVY)

add_textbox(slide, Inches(1), Inches(0.6), Inches(6), Inches(0.6),
            "FINANCE & REPORTING", font_size=14, color=GOLD, bold=True, font_name="Calibri")
add_accent_line(slide, Inches(1), Inches(1.1), Inches(1.5), GOLD)
add_textbox(slide, Inches(1), Inches(1.3), Inches(11), Inches(1),
            "Capital Governance with Full Visibility",
            font_size=36, color=WHITE, bold=True, font_name="Calibri Light")

fin_features = [
    ("Integration Hub", "Connect multiple external finance providers with encrypted credentials and pluggable architecture. Sync accounts and transactions with upsert semantics for idempotent reruns."),
    ("Health Monitoring", "Real-time sync status indicators (Healthy, Attention Needed, In Progress) with automated stale-sync detection and actionable error summaries."),
    ("Budgeting & Vendors", "Company-level budget tracking with vendor management, account mapping to projects, and 30-day transaction summaries at a glance."),
    ("PowerBI Reporting", "Embedded PowerBI dashboards for portfolio-wide business intelligence, with row-level security matching platform access controls."),
]

for i, (title, desc) in enumerate(fin_features):
    col = i % 2
    row = i // 2
    x = Inches(1) + col * Inches(5.8)
    y = Inches(2.6) + row * Inches(2.3)

    add_rounded_card(slide, x, y, Inches(5.4), Inches(2.0), DARK_BLUE, RGBColor(0x2A, 0x3A, 0x55))
    add_textbox(slide, x + Inches(0.4), y + Inches(0.25), Inches(4.6), Inches(0.4),
                title, font_size=18, color=GOLD, bold=True, font_name="Calibri")
    add_textbox(slide, x + Inches(0.4), y + Inches(0.7), Inches(4.6), Inches(1.1),
                desc, font_size=13, color=LIGHT_GRAY, font_name="Calibri")

# ============================================================
# SLIDE 11: PRODUCT SCREENSHOT - DATA ROOM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, NAVY)

add_textbox(slide, Inches(1), Inches(0.4), Inches(6), Inches(0.6),
            "PRODUCT SCREENSHOT", font_size=14, color=GOLD, bold=True, font_name="Calibri")
add_accent_line(slide, Inches(1), Inches(0.9), Inches(1.5), GOLD)
add_textbox(slide, Inches(1), Inches(1.1), Inches(11), Inches(0.7),
            "The Data Room in Action",
            font_size=32, color=WHITE, bold=True, font_name="Calibri Light")

screenshot_card = add_rounded_card(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(5.2),
                                    RGBColor(0x20, 0x30, 0x48), RGBColor(0x3B, 0x50, 0x70))

add_image_safe(slide, "screenshots/data-room-marketing-full.png",
               Inches(1.0), Inches(2.2), Inches(11.3))

# ============================================================
# SLIDE 12: SECURITY & ACCESS CONTROL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, NAVY)

add_textbox(slide, Inches(1), Inches(0.6), Inches(6), Inches(0.6),
            "SECURITY & ACCESS CONTROL", font_size=14, color=GOLD, bold=True, font_name="Calibri")
add_accent_line(slide, Inches(1), Inches(1.1), Inches(1.5), GOLD)
add_textbox(slide, Inches(1), Inches(1.3), Inches(11), Inches(1),
            "Enterprise-Grade Data Protection",
            font_size=36, color=WHITE, bold=True, font_name="Calibri Light")

security_items = [
    ("Multi-Company Access", "Granular authorization with a Canonical Effective-Access Resolver supporting read-only, contributor, and company admin roles across portfolios."),
    ("Module-Level Permissions", "Fine-grained permission enforcement at the module level -- control who can view, edit, or manage each area of the platform."),
    ("Role Profiles", "Detailed stakeholder definitions with Portfolio Hub Boundary Model for precise data visibility and operational access."),
    ("Encrypted Credentials", "All third-party integration credentials are encrypted at rest with environment-level isolation between tenants."),
]

for i, (title, desc) in enumerate(security_items):
    col = i % 2
    row = i // 2
    x = Inches(1) + col * Inches(5.8)
    y = Inches(2.6) + row * Inches(2.3)

    add_rounded_card(slide, x, y, Inches(5.4), Inches(2.0), DARK_BLUE, RGBColor(0x2A, 0x3A, 0x55))

    shield = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), y + Inches(0.35), Inches(0.5), Inches(0.5))
    shield.fill.solid()
    shield.fill.fore_color.rgb = GOLD
    shield.line.fill.background()
    add_textbox(slide, x + Inches(0.3), y + Inches(0.38), Inches(0.5), Inches(0.45),
                str(i+1), font_size=16, color=NAVY, bold=True, font_name="Calibri",
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, x + Inches(1.0), y + Inches(0.25), Inches(4.0), Inches(0.4),
                title, font_size=17, color=WHITE, bold=True, font_name="Calibri")
    add_textbox(slide, x + Inches(1.0), y + Inches(0.7), Inches(4.0), Inches(1.1),
                desc, font_size=12, color=LIGHT_GRAY, font_name="Calibri")

# ============================================================
# SLIDE 13: CLOSING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, NAVY)

add_image_safe(slide, "attached_assets/generated_images/real_estate_portfolio.png",
               Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)

overlay3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
overlay3.fill.solid()
overlay3.fill.fore_color.rgb = NAVY
set_shape_alpha(overlay3, 65)

add_accent_line(slide, Inches(4.2), Inches(2.2), Inches(5), GOLD)

add_textbox(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.5),
            "iliOS", font_size=72, color=WHITE, bold=True, font_name="Calibri Light",
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(4.0), Inches(10.3), Inches(0.8),
            "The Sun's Operating System", font_size=28, color=GOLD, bold=False, font_name="Calibri",
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.8),
            "AI-powered real estate investment management\nfor the modern portfolio.", font_size=18,
            color=LIGHT_GRAY, font_name="Calibri", alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.5), Inches(6.0), Inches(2.3), GOLD)

add_textbox(slide, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.5),
            "Schedule a Demo  |  info@ilios.energy", font_size=16, color=MID_GRAY,
            font_name="Calibri", alignment=PP_ALIGN.CENTER)

# ============================================================
# SAVE
# ============================================================
output_path = "screenshots/iliOS_Marketing_Deck_2026.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
